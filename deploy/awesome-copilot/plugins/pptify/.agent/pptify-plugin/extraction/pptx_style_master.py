from __future__ import annotations

import zipfile
from collections import Counter
from pathlib import Path
from typing import Any, Iterable
from xml.etree import ElementTree

EMU_PER_INCH = 914400
DRAWING_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


class PptxStyleMaster:
    """Extract compact design context from a reference PPTX for prompt-based generation."""

    def __init__(self, max_slides: int = 12, max_items: int = 10) -> None:
        self.max_slides = max_slides
        self.max_items = max_items

    def analyze(self, path: str | Path) -> dict[str, Any]:
        from pptx import Presentation

        pptx_path = Path(path)
        presentation = Presentation(str(pptx_path))
        slide_size = {
            "width": _inches(presentation.slide_width),
            "height": _inches(presentation.slide_height),
        }
        theme = _theme_from_package(pptx_path)

        colors: Counter[str] = Counter()
        fonts: Counter[str] = Counter()
        font_sizes: Counter[float] = Counter()
        shape_styles: Counter[str] = Counter()
        layout_names: Counter[str] = Counter()
        master_names: Counter[str] = Counter()
        slide_layouts: list[dict[str, Any]] = []

        _count_theme_tokens(theme, colors, fonts)
        for slide_index, slide in enumerate(presentation.slides, start=1):
            if slide_index > self.max_slides:
                break
            slide_context = _slide_design_context(slide, slide_index, slide_size, self.max_items)
            slide_layouts.append(slide_context)
            layout_names[slide_context["template_layout"]] += 1
            master_names[slide_context["template_master"]] += 1
            colors.update(slide_context.pop("_colors"))
            fonts.update(slide_context.pop("_fonts"))
            font_sizes.update(slide_context.pop("_font_sizes"))
            shape_styles.update(slide_context.pop("_shape_styles"))

        return {
            "styles": {
                "colors": _top_items(colors, self.max_items),
                "fonts": _top_items(fonts, self.max_items),
                "font_sizes": _top_items(font_sizes, self.max_items),
                "shape_styles": _top_items(shape_styles, self.max_items),
            },
            "brands": _brand_context(colors, fonts, theme, self.max_items),
            "template": _template_context(presentation, slide_size, theme, layout_names, master_names, self.max_items),
            "layout": {
                "analyzed_slide_count": len(slide_layouts),
                "layout_usage": _top_items(layout_names, self.max_items),
                "master_usage": _top_items(master_names, self.max_items),
                "slides": slide_layouts,
            },
        }


def extract_pptx_style_master(path: str | Path, max_slides: int = 12, max_items: int = 10) -> dict[str, Any]:
    return PptxStyleMaster(max_slides=max_slides, max_items=max_items).analyze(path)


def _slide_design_context(slide, slide_index: int, slide_size: dict[str, float], max_items: int) -> dict[str, Any]:
    colors: Counter[str] = Counter()
    fonts: Counter[str] = Counter()
    font_sizes: Counter[float] = Counter()
    shape_styles: Counter[str] = Counter()
    object_counts: Counter[str] = Counter()
    regions: Counter[str] = Counter()
    placeholders: list[dict[str, Any]] = []
    object_samples: list[dict[str, Any]] = []
    boxes: list[dict[str, float]] = []

    for shape in _iter_shapes(slide.shapes):
        kind = _shape_kind(shape)
        bbox = _bbox(shape)
        boxes.append(bbox)
        object_counts[kind] += 1
        regions[_region(bbox, slide_size)] += 1

        shape_colors = _shape_colors(shape)
        colors.update(shape_colors.values())
        shape_text = _text_preview(shape)
        text_styles = _text_styles(shape)
        fonts.update(text_styles["fonts"])
        font_sizes.update(text_styles["font_sizes"])
        colors.update(text_styles["colors"])

        style_signature = _style_signature(shape_colors, text_styles)
        if style_signature:
            shape_styles[style_signature] += 1

        if getattr(shape, "is_placeholder", False) and len(placeholders) < max_items:
            placeholders.append(_placeholder_context(shape, bbox))

        if len(object_samples) < max_items:
            sample: dict[str, Any] = {
                "kind": kind,
                "role": _shape_role(shape, kind),
                "bbox": bbox,
                "region": _region(bbox, slide_size),
            }
            if shape_text:
                sample["text"] = shape_text
            if shape_colors:
                sample["colors"] = shape_colors
            if text_styles["fonts"]:
                sample["fonts"] = _top_items(text_styles["fonts"], 3)
            if text_styles["font_sizes"]:
                sample["font_sizes"] = _top_items(text_styles["font_sizes"], 3)
            object_samples.append(sample)

    return {
        "index": slide_index,
        "template_layout": _slide_layout_name(slide),
        "template_master": _slide_master_name(slide),
        "object_counts": dict(sorted(object_counts.items())),
        "placeholder_count": len(placeholders),
        "placeholders": placeholders,
        "dominant_regions": _top_items(regions, max_items),
        "dominant_flow": _dominant_flow(boxes, slide_size),
        "occupied_area_ratio": _occupied_area_ratio(boxes, slide_size),
        "objects": object_samples,
        "_colors": colors,
        "_fonts": fonts,
        "_font_sizes": font_sizes,
        "_shape_styles": shape_styles,
    }


def _template_context(
    presentation,
    slide_size: dict[str, float],
    theme: dict[str, Any],
    layout_names: Counter[str],
    master_names: Counter[str],
    max_items: int,
) -> dict[str, Any]:
    masters: list[dict[str, Any]] = []
    try:
        for master_index, master in enumerate(presentation.slide_masters, start=1):
            masters.append(
                {
                    "index": master_index,
                    "name": str(getattr(master, "name", f"Master {master_index}") or f"Master {master_index}"),
                    "layout_count": len(master.slide_layouts),
                }
            )
    except (AttributeError, TypeError):
        masters = []

    return {
        "slide_size": slide_size,
        "theme": theme,
        "masters": masters[:max_items],
        "layout_usage": _top_items(layout_names, max_items),
        "master_usage": _top_items(master_names, max_items),
    }


def _brand_context(colors: Counter[str], fonts: Counter[str], theme: dict[str, Any], max_items: int) -> dict[str, Any]:
    theme_colors = theme.get("colors", {}) if isinstance(theme.get("colors"), dict) else {}
    theme_accents = [value for name, value in theme_colors.items() if str(name).startswith("accent")]
    palette = _ranked_colors(colors, include_neutral=False)
    if not palette:
        palette = [color for color in theme_accents if _is_hex_color(color)]
    neutral_palette = _ranked_colors(colors, include_neutral=True, only_neutral=True)
    font_values = [str(item["value"]) for item in _top_items(fonts, max_items)]
    primary_color = palette[0] if palette else None
    accent_colors = _dedupe([*palette, *theme_accents])[:max_items]

    return {
        "theme_name": theme.get("name"),
        "primary_color": primary_color,
        "accent_colors": accent_colors,
        "neutral_colors": neutral_palette[:max_items],
        "fonts": font_values[:max_items],
        "theme_colors": theme_colors,
        "theme_fonts": theme.get("fonts", {}),
    }


def _theme_from_package(path: Path) -> dict[str, Any]:
    theme_paths: list[str]
    try:
        with zipfile.ZipFile(path) as package:
            theme_paths = sorted(name for name in package.namelist() if name.startswith("ppt/theme/theme") and name.endswith(".xml"))
            if not theme_paths:
                return {"name": None, "colors": {}, "fonts": {}}
            root = ElementTree.fromstring(package.read(theme_paths[0]))
    except (zipfile.BadZipFile, KeyError, ElementTree.ParseError):
        return {"name": None, "colors": {}, "fonts": {}}

    theme = {
        "name": root.attrib.get("name"),
        "path": theme_paths[0],
        "colors": {},
        "fonts": {},
    }
    color_scheme = root.find(f".//{DRAWING_NS}clrScheme")
    if color_scheme is not None:
        colors: dict[str, str] = {}
        for color_node in list(color_scheme):
            color_value = _theme_color_value(color_node)
            if color_value:
                colors[color_node.tag.rsplit("}", 1)[-1]] = color_value
        theme["colors"] = colors

    font_scheme = root.find(f".//{DRAWING_NS}fontScheme")
    if font_scheme is not None:
        fonts: dict[str, str] = {}
        for key, node_name in (("major", "majorFont"), ("minor", "minorFont")):
            latin = font_scheme.find(f".//{DRAWING_NS}{node_name}/{DRAWING_NS}latin")
            if latin is not None and latin.attrib.get("typeface"):
                fonts[key] = latin.attrib["typeface"]
        theme["fonts"] = fonts
    return theme


def _theme_color_value(color_node: ElementTree.Element) -> str | None:
    srgb = color_node.find(f".//{DRAWING_NS}srgbClr")
    if srgb is not None and srgb.attrib.get("val"):
        return _normalize_hex(srgb.attrib["val"])
    system = color_node.find(f".//{DRAWING_NS}sysClr")
    if system is not None and system.attrib.get("lastClr"):
        return _normalize_hex(system.attrib["lastClr"])
    return None


def _count_theme_tokens(theme: dict[str, Any], colors: Counter[str], fonts: Counter[str]) -> None:
    for color in theme.get("colors", {}).values() if isinstance(theme.get("colors"), dict) else []:
        if _is_hex_color(color):
            colors[color] += 1
    for font in theme.get("fonts", {}).values() if isinstance(theme.get("fonts"), dict) else []:
        if font:
            fonts[str(font)] += 1


def _iter_shapes(shapes) -> Iterable[Any]:
    for shape in shapes:
        yield shape
        if hasattr(shape, "shapes"):
            yield from _iter_shapes(shape.shapes)


def _shape_kind(shape) -> str:
    shape_type = str(getattr(getattr(shape, "shape_type", "unknown"), "name", "unknown")).lower()
    if getattr(shape, "has_table", False):
        return "table"
    if getattr(shape, "has_chart", False):
        return "chart"
    if "picture" in shape_type or _has_image(shape):
        return "image"
    if "line" in shape_type or "connector" in shape_type or "freeform" in shape_type:
        return "line"
    if getattr(shape, "has_text_frame", False) and getattr(shape, "text", "").strip():
        return "text"
    if hasattr(shape, "shapes"):
        return "group"
    return "shape"


def _has_image(shape) -> bool:
    try:
        return getattr(shape, "image", None) is not None
    except (AttributeError, TypeError, ValueError):
        return False


def _shape_role(shape, kind: str) -> str:
    if getattr(shape, "is_placeholder", False):
        try:
            return str(shape.placeholder_format.type).split(".")[-1].lower()
        except (AttributeError, ValueError):
            return "placeholder"
    name = str(getattr(shape, "name", "") or "").strip().lower()
    if "title" in name:
        return "title"
    return kind


def _shape_colors(shape) -> dict[str, str]:
    colors: dict[str, str] = {}
    fill = _format_color(_safe_attr(_safe_attr(shape, "fill"), "fore_color"))
    if fill:
        colors["fill"] = fill
    line = _format_color(_safe_attr(_safe_attr(shape, "line"), "color"))
    if line:
        colors["line"] = line
    return colors


def _safe_attr(value: Any, name: str) -> Any:
    if value is None:
        return None
    try:
        return getattr(value, name)
    except (AttributeError, TypeError, ValueError):
        return None


def _text_styles(shape) -> dict[str, Counter[Any]]:
    fonts: Counter[str] = Counter()
    font_sizes: Counter[float] = Counter()
    colors: Counter[str] = Counter()
    text_frame = getattr(shape, "text_frame", None)
    if text_frame is None:
        return {"fonts": fonts, "font_sizes": font_sizes, "colors": colors}

    for paragraph in text_frame.paragraphs:
        _count_font(paragraph.font, fonts, font_sizes, colors)
        for run in paragraph.runs:
            _count_font(run.font, fonts, font_sizes, colors)
    return {"fonts": fonts, "font_sizes": font_sizes, "colors": colors}


def _count_font(font, fonts: Counter[str], font_sizes: Counter[float], colors: Counter[str]) -> None:
    name = getattr(font, "name", None)
    if name:
        fonts[str(name)] += 1
    size = getattr(font, "size", None)
    if size is not None:
        font_sizes[round(size.pt, 2)] += 1
    color = _format_color(getattr(font, "color", None))
    if color:
        colors[color] += 1


def _format_color(color_format) -> str | None:
    if color_format is None:
        return None
    try:
        rgb = color_format.rgb
    except (AttributeError, TypeError, ValueError):
        rgb = None
    if rgb is not None:
        return _normalize_hex(str(rgb))
    try:
        theme_color = color_format.theme_color
    except (AttributeError, TypeError, ValueError):
        theme_color = None
    if theme_color:
        token = str(theme_color).split(".")[-1].lower()
        return f"theme:{token}"
    return None


def _style_signature(shape_colors: dict[str, str], text_styles: dict[str, Counter[Any]]) -> str:
    parts: list[str] = []
    if shape_colors.get("fill"):
        parts.append(f"fill={shape_colors['fill']}")
    if shape_colors.get("line"):
        parts.append(f"line={shape_colors['line']}")
    font = _top_value(text_styles["fonts"])
    if font:
        parts.append(f"font={font}")
    font_size = _top_value(text_styles["font_sizes"])
    if font_size:
        parts.append(f"font_size={font_size}")
    return "; ".join(parts)


def _placeholder_context(shape, bbox: dict[str, float]) -> dict[str, Any]:
    context: dict[str, Any] = {
        "name": str(getattr(shape, "name", "") or ""),
        "bbox": bbox,
    }
    try:
        context["type"] = str(shape.placeholder_format.type).split(".")[-1].lower()
        context["idx"] = int(shape.placeholder_format.idx)
    except (AttributeError, ValueError):
        context["type"] = "placeholder"
    return context


def _text_preview(shape, max_chars: int = 120) -> str:
    text = " ".join(str(getattr(shape, "text", "")).split())
    return text[:max_chars]


def _bbox(shape) -> dict[str, float]:
    return {
        "x": _inches(getattr(shape, "left", 0)),
        "y": _inches(getattr(shape, "top", 0)),
        "width": max(0.0, _inches(getattr(shape, "width", 0))),
        "height": max(0.0, _inches(getattr(shape, "height", 0))),
    }


def _region(bbox: dict[str, float], slide_size: dict[str, float]) -> str:
    width = max(slide_size["width"], 0.01)
    height = max(slide_size["height"], 0.01)
    center_x = (bbox["x"] + bbox["width"] / 2) / width
    center_y = (bbox["y"] + bbox["height"] / 2) / height
    horizontal = "left" if center_x < 0.34 else "right" if center_x > 0.66 else "center"
    vertical = "top" if center_y < 0.34 else "bottom" if center_y > 0.66 else "middle"
    return f"{vertical}_{horizontal}"


def _dominant_flow(boxes: list[dict[str, float]], slide_size: dict[str, float]) -> str:
    if len(boxes) < 2:
        return "single"
    centers_x = [(box["x"] + box["width"] / 2) / max(slide_size["width"], 0.01) for box in boxes]
    centers_y = [(box["y"] + box["height"] / 2) / max(slide_size["height"], 0.01) for box in boxes]
    spread_x = max(centers_x) - min(centers_x)
    spread_y = max(centers_y) - min(centers_y)
    if len(boxes) >= 4 and spread_x > 0.32 and spread_y > 0.32:
        return "grid"
    if spread_x > 0.42 and spread_y > 0.42:
        return "grid"
    if len(boxes) >= 3 and spread_x > 0.42:
        return "grid"
    if spread_x > spread_y * 1.4:
        return "row"
    if spread_y > spread_x * 1.4:
        return "column"
    return "overlay_or_balanced"


def _occupied_area_ratio(boxes: list[dict[str, float]], slide_size: dict[str, float]) -> float:
    slide_area = max(slide_size["width"] * slide_size["height"], 0.01)
    object_area = sum(box["width"] * box["height"] for box in boxes)
    return round(min(object_area / slide_area, 1.0), 3)


def _slide_layout_name(slide) -> str:
    try:
        return str(slide.slide_layout.name or "unnamed_layout")
    except AttributeError:
        return "unknown_layout"


def _slide_master_name(slide) -> str:
    try:
        master = slide.slide_layout.slide_master
        return str(master.name or "unnamed_master")
    except AttributeError:
        return "unknown_master"


def _top_items(counter: Counter[Any], limit: int) -> list[dict[str, Any]]:
    return [{"value": value, "count": count} for value, count in counter.most_common(limit)]


def _top_value(counter: Counter[Any]) -> Any | None:
    if not counter:
        return None
    return counter.most_common(1)[0][0]


def _ranked_colors(colors: Counter[str], include_neutral: bool, only_neutral: bool = False) -> list[str]:
    ranked: list[str] = []
    for color, _count in colors.most_common():
        if not _is_hex_color(color):
            continue
        neutral = _is_neutral(color)
        if only_neutral and not neutral:
            continue
        if not include_neutral and neutral:
            continue
        ranked.append(color)
    return ranked


def _is_neutral(color: str) -> bool:
    if not _is_hex_color(color):
        return False
    red = int(color[1:3], 16)
    green = int(color[3:5], 16)
    blue = int(color[5:7], 16)
    return max(red, green, blue) - min(red, green, blue) <= 18


def _is_hex_color(value: Any) -> bool:
    return isinstance(value, str) and len(value) == 7 and value.startswith("#")


def _normalize_hex(value: str) -> str:
    stripped = value.strip().lstrip("#")
    if len(stripped) >= 6:
        return f"#{stripped[:6].upper()}"
    return f"#{stripped.upper()}"


def _dedupe(values: Iterable[str]) -> list[str]:
    deduped: list[str] = []
    for value in values:
        if value and value not in deduped:
            deduped.append(value)
    return deduped


def _inches(value: int) -> float:
    return round(int(value) / EMU_PER_INCH, 4)
