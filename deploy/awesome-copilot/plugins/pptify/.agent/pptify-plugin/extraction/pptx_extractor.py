from __future__ import annotations

import json
import posixpath
import sys
import zipfile
from base64 import b64encode
from collections import Counter
from pathlib import Path
from typing import Any
from xml.etree import ElementTree

# Allow importing sibling pptx_style_master when run as a standalone script
sys.path.insert(0, str(Path(__file__).parent))
from pptx_style_master import PptxStyleMaster

EMU_PER_INCH = 914400
DRAWING_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


class PptxExtractor:
    def prompt_context(self, path: str | Path, max_chars: int = 16000) -> dict[str, Any]:
        from pptx import Presentation

        pptx_path = Path(path)
        presentation = Presentation(str(pptx_path))
        style_context = PptxStyleMaster().analyze(pptx_path)
        slides: list[dict[str, Any]] = []
        used_chars = 0
        for slide_index, slide in enumerate(presentation.slides, start=1):
            texts = _slide_text_fragments(slide.shapes)
            trimmed_texts: list[str] = []
            for text in texts:
                if used_chars >= max_chars:
                    break
                cleaned = _compact_text(text)
                if not cleaned:
                    continue
                remaining = max_chars - used_chars
                clipped = cleaned[: min(500, remaining)]
                trimmed_texts.append(clipped)
                used_chars += len(clipped)
            title = trimmed_texts[0] if trimmed_texts else f"Slide {slide_index}"
            slides.append(
                {
                    "index": slide_index,
                    "title": title[:120],
                    "text": trimmed_texts[:12],
                    "shape_count": len(slide.shapes),
                }
            )
        media_files, embedded_files = _package_asset_counts(pptx_path)
        return {
            "source": str(pptx_path),
            "slide_count": len(slides),
            "slide_size": {
                "width": _inches(presentation.slide_width),
                "height": _inches(presentation.slide_height),
            },
            "package_media_files": media_files,
            "embedded_files": embedded_files,
            "styles": style_context["styles"],
            "brands": style_context["brands"],
            "template": style_context["template"],
            "layout": style_context["layout"],
            "slides": slides,
        }

    def extract_file(self, path: str | Path, output_dir: str | Path | None = None, extract_media: bool = True) -> dict[str, Any]:
        from pptx import Presentation

        pptx_path = Path(path)
        presentation = Presentation(str(pptx_path))
        asset_dir = None
        embed_media = extract_media and output_dir is None
        if output_dir and extract_media:
            asset_dir = Path(output_dir) / f"{pptx_path.stem}_assets"
            asset_dir.mkdir(parents=True, exist_ok=True)
            _extract_package_media(pptx_path, asset_dir)

        notes_by_slide = _notes_by_slide(pptx_path)
        slides: list[dict[str, Any]] = []
        stats: Counter[str] = Counter()
        max_shapes = 0
        max_nested = 0
        for slide_index, slide in enumerate(presentation.slides, start=1):
            tree, slide_stats, render_elements = self._extract_slide(
                slide=slide,
                slide_index=slide_index,
                slide_size=(_inches(presentation.slide_width), _inches(presentation.slide_height)),
                source_path=pptx_path,
                asset_dir=asset_dir,
                embed_media=embed_media,
                notes=notes_by_slide.get(slide_index, []),
            )
            stats.update(slide_stats)
            max_shapes = max(max_shapes, slide_stats["top_level_shapes"])
            max_nested = max(max_nested, slide_stats["nested_shapes"])
            slides.append(
                {
                    "id": tree["id"],
                    "title": tree["title"],
                    "slide_size": tree["slide_size"],
                    "preserve_coordinates": True,
                    "render_mode": "ooxml",
                    "ooxml_elements": render_elements,
                    "layout_tree": tree,
                }
            )

        media_files, embedded_files = _package_asset_counts(pptx_path)
        style_context = PptxStyleMaster().analyze(pptx_path)
        summary = {
            "source": str(pptx_path),
            "slide_count": len(slides),
            "slide_size": {
                "width": _inches(presentation.slide_width),
                "height": _inches(presentation.slide_height),
            },
            "top_level_shapes": int(stats["top_level_shapes"]),
            "nested_shapes": int(stats["nested_shapes"]),
            "max_shapes_on_slide": max_shapes,
            "max_nested_shapes_on_slide": max_nested,
            "groups": int(stats["groups"]),
            "tables": int(stats["tables"]),
            "charts": int(stats["charts"]),
            "images": int(stats["images"]),
            "text_objects": int(stats["text_objects"]),
            "placeholders": int(stats["placeholders"]),
            "lines_or_freeforms": int(stats["lines_or_freeforms"]),
            "non_ascii_text": bool(stats["non_ascii_text"]),
            "notes_slides": int(stats["notes_slides"]),
            "package_media_files": media_files,
            "embedded_files": embedded_files,
            "styles": style_context["styles"],
            "brands": style_context["brands"],
            "template": style_context["template"],
            "layout": style_context["layout"],
        }
        return {
            "source_pptx": str(pptx_path.resolve()),
            "render_mode": "ooxml",
            "summary": summary,
            "slides": slides,
        }

    def extract_path(self, path: str | Path, output_dir: str | Path, extract_media: bool = True) -> dict[str, Any]:
        source = Path(path)
        output = Path(output_dir)
        output.mkdir(parents=True, exist_ok=True)
        files = sorted(source.glob("*.pptx")) if source.is_dir() else [source]
        decks = []
        for pptx_file in files:
            deck = self.extract_file(pptx_file, output, extract_media=extract_media)
            json_path = output / f"{pptx_file.stem}.pptify.json"
            json_path.write_text(json.dumps(deck, indent=2, ensure_ascii=False), encoding="utf-8")
            decks.append({"pptx": str(pptx_file), "json": str(json_path), "summary": deck["summary"]})
        manifest = {"source": str(source), "decks": decks}
        (output / "manifest.json").write_text(json.dumps(manifest, indent=2, ensure_ascii=False), encoding="utf-8")
        return manifest

    def analyze_path(self, path: str | Path) -> dict[str, Any]:
        source = Path(path)
        files = sorted(source.glob("*.pptx")) if source.is_dir() else [source]
        return {"source": str(source), "decks": [self.extract_file(file, extract_media=False)["summary"] for file in files]}

    def _extract_slide(
        self,
        slide,
        slide_index: int,
        slide_size: tuple[float, float],
        source_path: Path,
        asset_dir: Path | None,
        embed_media: bool,
        notes: list[str],
    ) -> tuple[dict[str, Any], Counter[str], list[dict[str, Any]]]:
        root_id = f"slide_{slide_index}_root"
        root_group: dict[str, Any] = {
            "id": root_id,
            "role": "slide",
            "layout_mode": "absolute",
            "object_ids": [],
            "group_ids": [],
            "constraints": {},
            "collision_policy": "relaxed",
            "bbox": {"x": 0, "y": 0, "width": slide_size[0], "height": slide_size[1]},
        }
        groups: dict[str, dict[str, Any]] = {root_id: root_group}
        objects: dict[str, dict[str, Any]] = {}
        stats: Counter[str] = Counter(top_level_shapes=len(slide.shapes), notes_slides=1 if notes else 0)
        z_index = 0

        def walk(shapes, parent_group_id: str, prefix: str) -> None:
            nonlocal z_index
            for shape_index, shape in enumerate(shapes, start=1):
                z_index += 1
                stats["nested_shapes"] += 1
                shape_type = _shape_type_name(shape)
                if _is_group(shape):
                    group_id = f"{prefix}_group_{shape_index}"
                    groups[parent_group_id]["group_ids"].append(group_id)
                    groups[group_id] = {
                        "id": group_id,
                        "role": "extracted_group",
                        "layout_mode": "absolute",
                        "object_ids": [],
                        "group_ids": [],
                        "constraints": {},
                        "collision_policy": "relaxed",
                        "bbox": _bbox(shape),
                    }
                    stats["groups"] += 1
                    walk(shape.shapes, group_id, group_id)
                    continue

                object_id = f"{prefix}_shape_{shape_index}"
                slide_object = self._extract_object(shape, object_id, z_index, shape_type, source_path, asset_dir)
                objects[object_id] = slide_object
                groups[parent_group_id]["object_ids"].append(object_id)
                stats[_stat_key(slide_object)] += 1
                if getattr(shape, "is_placeholder", False):
                    stats["placeholders"] += 1
                if _contains_non_ascii(slide_object["content"].get("text", "")):
                    stats["non_ascii_text"] += 1

        walk(slide.shapes, root_id, f"slide_{slide_index}")
        render_elements = [
            _render_element(shape, f"slide_{slide_index}_element_{element_index}", source_path, asset_dir, embed_media)
            for element_index, shape in enumerate(slide.shapes, start=1)
        ]
        title = _slide_title(objects.values()) or f"Slide {slide_index}"
        tree: dict[str, Any] = {
            "id": f"slide_{slide_index}",
            "title": title,
            "slide_size": {"width": slide_size[0], "height": slide_size[1]},
            "root_group_id": root_id,
            "groups": groups,
            "objects": objects,
            "notes": notes,
        }
        return tree, stats, render_elements

    def _extract_object(self, shape, object_id: str, z_index: int, shape_type: str, source_path: Path, asset_dir: Path | None) -> dict[str, Any]:
        kind = _kind(shape, shape_type)
        content: dict[str, Any] = {"source_shape_type": shape_type}
        style: dict[str, Any] = {}
        if kind == "text":
            content["text"] = getattr(shape, "text", "")
            style.update(_text_style(shape))
        elif kind == "table":
            content["rows"] = [[cell.text for cell in row.cells] for row in shape.table.rows]
            style["font_size"] = 8
        elif kind == "image":
            content["alt"] = getattr(shape, "name", "image")
            if asset_dir is not None:
                image_data = _image_data(shape)
                if image_data is None:
                    content["missing_embedded_image"] = True
                else:
                    blob, extension, relationship_id, content_type = image_data
                    asset_path = asset_dir / f"{source_path.stem}_{object_id}.{extension}"
                    asset_path.write_bytes(blob)
                    content["path"] = str(asset_path)
                    content["content_type"] = content_type
                    if relationship_id:
                        content["media_relationship_id"] = relationship_id
        elif kind == "chart":
            content["title"] = getattr(shape, "name", "chart")
        elif kind == "line":
            box = _bbox(shape)
            x, y, w, h = box["x"], box["y"], box["width"], box["height"]
            content.update({"x1": x, "y1": y, "x2": x + w, "y2": y + h})
            style["line"] = "#6B7280"
        elif getattr(shape, "has_text_frame", False) and getattr(shape, "text", ""):
            content["text"] = shape.text
            style.update(_text_style(shape))
        else:
            content["shape"] = "rect"

        classification = "content" if kind in {"text", "table", "image", "chart"} else "layout_design"
        if kind == "shape" and content.get("text"):
            classification = "content"
        return {
            "id": object_id,
            "kind": kind,
            "role": _role(shape, kind),
            "classification": classification,
            "content": content,
            "style": style,
            "constraints": {"source_name": getattr(shape, "name", "")},
            "bbox": _bbox(shape),
            "z_index": z_index,
        }


def _inches(value: int) -> float:
    return round(int(value) / EMU_PER_INCH, 4)


def _bbox(shape) -> dict[str, float]:
    return {
        "x": _inches(getattr(shape, "left", 0) or 0),
        "y": _inches(getattr(shape, "top", 0) or 0),
        "width": max(0.0, _inches(getattr(shape, "width", 0) or 0)),
        "height": max(0.0, _inches(getattr(shape, "height", 0) or 0)),
    }


def _shape_type_name(shape) -> str:
    shape_type = getattr(shape, "shape_type", "unknown")
    return str(getattr(shape_type, "name", shape_type)).lower()


def _is_group(shape) -> bool:
    return hasattr(shape, "shapes") and "group" in _shape_type_name(shape)


def _kind(shape, shape_type: str) -> str:
    if getattr(shape, "has_table", False):
        return "table"
    if getattr(shape, "has_chart", False):
        return "chart"
    if "picture" in shape_type or hasattr(shape, "image"):
        return "image"
    if "line" in shape_type or "freeform" in shape_type or "connector" in shape_type:
        return "line"
    if getattr(shape, "has_text_frame", False) and getattr(shape, "text", "").strip():
        return "text"
    return "shape"


def _role(shape, kind: str) -> str:
    if getattr(shape, "is_placeholder", False):
        return "placeholder"
    if kind == "text":
        return "text"
    return kind


def _text_style(shape) -> dict[str, Any]:
    style: dict[str, Any] = {"font_size": 12}
    try:
        paragraph = shape.text_frame.paragraphs[0]
        run = paragraph.runs[0] if paragraph.runs else None
        font = run.font if run is not None else paragraph.font
        if font.size is not None:
            style["font_size"] = round(font.size.pt, 2)
        if font.bold is not None:
            style["bold"] = bool(font.bold)
        if font.name:
            style["font"] = font.name
    except (AttributeError, IndexError):
        pass
    return style


def _image_data(shape) -> tuple[bytes, str, str | None, str] | None:
    try:
        image = shape.image
    except ValueError:
        image = None
    if image is not None:
        return image.blob, str(image.ext), None, image.content_type

    for relationship_id in _embedded_relationship_ids(shape):
        try:
            part = shape.part.related_part(relationship_id)
        except KeyError:
            continue
        blob = getattr(part, "blob", None)
        if not blob:
            continue
        extension = Path(str(getattr(part, "partname", "media.bin"))).suffix.lstrip(".") or _extension_from_content_type(
            getattr(part, "content_type", "")
        )
        return blob, extension or "bin", relationship_id, getattr(part, "content_type", "")
    return None


def _render_element(shape, element_id: str, source_path: Path, asset_dir: Path | None, embed_media: bool) -> dict[str, Any]:
    return {
        "id": element_id,
        "xml": shape._element.xml,
        "relationships": _relationship_payloads(shape, element_id, source_path, asset_dir, embed_media),
    }


def _relationship_payloads(
    shape,
    element_id: str,
    source_path: Path,
    asset_dir: Path | None,
    embed_media: bool,
) -> list[dict[str, Any]]:
    payloads: list[dict[str, Any]] = []
    for relationship_id in _embedded_relationship_ids(shape):
        try:
            relationship = shape.part.rels[relationship_id]
        except KeyError:
            continue
        payload: dict[str, Any] = {
            "source_rid": relationship_id,
            "reltype": relationship.reltype,
            "target_ref": relationship.target_ref,
            "is_external": bool(relationship.is_external),
        }
        if relationship.is_external:
            payloads.append(payload)
            continue
        part = relationship.target_part
        blob = getattr(part, "blob", None)
        if blob:
            content_type = getattr(part, "content_type", "")
            extension = Path(str(getattr(part, "partname", "media.bin"))).suffix.lstrip(".") or _extension_from_content_type(content_type)
            payload.update({"content_type": content_type, "extension": extension or "bin"})
            if asset_dir is not None:
                asset_path = asset_dir / f"{source_path.stem}_{element_id}_{relationship_id}.{payload['extension']}"
                asset_path.write_bytes(blob)
                payload["path"] = str(asset_path)
            elif embed_media:
                payload["blob_base64"] = b64encode(blob).decode("ascii")
        payloads.append(payload)
    return payloads


def _embedded_relationship_ids(shape) -> list[str]:
    relationship_ids: list[str] = []
    try:
        nodes = shape._element.iter()
    except AttributeError:
        return relationship_ids
    for node in nodes:
        for attribute_name, value in node.attrib.items():
            if attribute_name.endswith("}embed") and value not in relationship_ids:
                relationship_ids.append(value)
    return relationship_ids


def _extension_from_content_type(content_type: str) -> str:
    mapping = {
        "image/svg+xml": "svg",
        "image/png": "png",
        "image/jpeg": "jpg",
        "image/jpg": "jpg",
        "image/gif": "gif",
        "image/bmp": "bmp",
        "image/x-emf": "emf",
        "image/x-wmf": "wmf",
    }
    return mapping.get(content_type.lower(), "bin")


def _stat_key(slide_object: dict[str, Any]) -> str:
    kind = slide_object["kind"]
    if kind == "table":
        return "tables"
    if kind == "chart":
        return "charts"
    if kind == "image":
        return "images"
    if kind == "text":
        return "text_objects"
    if kind == "line":
        return "lines_or_freeforms"
    return "shapes"


def _slide_title(objects) -> str:
    for slide_object in objects:
        text = str(slide_object["content"].get("text", "")).strip()
        if text:
            return text.splitlines()[0][:100]
    return ""


def _contains_non_ascii(value: str) -> bool:
    return any(ord(character) > 127 for character in value)


def _slide_text_fragments(shapes) -> list[str]:
    fragments: list[str] = []
    for shape in shapes:
        if hasattr(shape, "shapes"):
            fragments.extend(_slide_text_fragments(shape.shapes))
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                row_text = " | ".join(_compact_text(cell.text) for cell in row.cells if _compact_text(cell.text))
                if row_text:
                    fragments.append(row_text)
        text = getattr(shape, "text", "")
        if text:
            fragments.append(text)
    deduped: list[str] = []
    seen: set[str] = set()
    for fragment in fragments:
        cleaned = _compact_text(fragment)
        if cleaned and cleaned not in seen:
            seen.add(cleaned)
            deduped.append(cleaned)
    return deduped


def _compact_text(value: str) -> str:
    return " ".join(str(value).split())


def _package_asset_counts(path: Path) -> tuple[int, int]:
    with zipfile.ZipFile(path) as package:
        names = package.namelist()
    media = sum(1 for name in names if name.startswith("ppt/media/"))
    embedded = sum(1 for name in names if name.startswith("ppt/embeddings/"))
    return media, embedded


def _extract_package_media(path: Path, asset_dir: Path) -> None:
    with zipfile.ZipFile(path) as package:
        for name in package.namelist():
            if not name.startswith("ppt/media/"):
                continue
            target = asset_dir / Path(name).name
            target.write_bytes(package.read(name))


def _notes_by_slide(path: Path) -> dict[int, list[str]]:
    notes: dict[int, list[str]] = {}
    with zipfile.ZipFile(path) as package:
        names = set(package.namelist())
        slide_rels = sorted(name for name in names if name.startswith("ppt/slides/_rels/slide") and name.endswith(".xml.rels"))
        for rel_path in slide_rels:
            slide_number = int(Path(rel_path).name.removeprefix("slide").removesuffix(".xml.rels"))
            rel_root = ElementTree.fromstring(package.read(rel_path))
            for rel in rel_root:
                if "notesSlide" not in rel.attrib.get("Type", ""):
                    continue
                target = rel.attrib.get("Target", "")
                notes_path = posixpath.normpath(posixpath.join("ppt/slides", target))
                if notes_path not in names:
                    notes_path = posixpath.normpath(posixpath.join("ppt/slides/_rels", target))
                if notes_path not in names:
                    continue
                notes_root = ElementTree.fromstring(package.read(notes_path))
                text = "\n".join(node.text for node in notes_root.iter(f"{DRAWING_NS}t") if node.text)
                if text.strip():
                    notes[slide_number] = [text.strip()]
    return notes
